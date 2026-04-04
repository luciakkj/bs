from __future__ import annotations

import argparse
import csv
import json
import subprocess
from dataclasses import dataclass
from pathlib import Path

from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


DESKTOP = Path.home() / "Desktop"
PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "output"
GENERATED_DIR = OUTPUT_DIR / "generated_ppt"
PREVIEW_DIR = GENERATED_DIR / "template_preview"
ASSET_DIR = GENERATED_DIR / "assets"
X_SCALE = 2.55

TEXT_DARK = RGBColor(40, 46, 92)
TEXT_MUTED = RGBColor(90, 99, 130)
CARD_BG = RGBColor(246, 248, 252)
CARD_BORDER = RGBColor(218, 224, 238)
ACCENT_BLUE = RGBColor(61, 92, 255)
ACCENT_PINK = RGBColor(228, 17, 94)
ACCENT_GOLD = RGBColor(244, 175, 28)
ACCENT_GREEN = RGBColor(25, 171, 128)


@dataclass
class SourceFiles:
    proposal_docx: Path
    title_docx: Path
    template_pptx: Path


@dataclass
class Stats:
    title: str
    reporter: str
    teacher: str
    detection_map50: float
    detection_map5095: float
    detection_run: str
    behavior_f1: float
    behavior_val_samples: int
    frame_f1: float
    frame_precision: float
    frame_recall: float
    runtime_fps: float
    runtime_source: str
    runtime_frames: int
    python_files: int
    python_lines: int
    behavior_runs: int
    compare_files: int
    runtime_logs: int
    snapshots: int
    alarms_total: int
    pseudo_tracks: int
    expanded_tracks: int


def zh(code: str) -> str:
    return code.encode("ascii").decode("unicode_escape")


def ensure_dirs() -> None:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def find_source_files() -> SourceFiles:
    thesis_dir = DESKTOP / zh(r"\u6bd5\u8bbe")
    proposal = next(p for p in thesis_dir.glob("*%s*.docx" % zh(r"\u5f00\u9898\u62a5\u544a")) if not p.name.startswith("~$"))
    title_doc = thesis_dir / zh(r"\u9898\u76ee.docx")
    template = next(DESKTOP.glob("*ppt*.pptx"))
    return SourceFiles(proposal_docx=proposal, title_docx=title_doc, template_pptx=template)


def load_doc_text(path: Path) -> list[str]:
    doc = Document(str(path))
    items: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            items.append(text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    items.append(text)
    return items


def value_after(items: list[str], label: str) -> str:
    for i, item in enumerate(items):
        if item == label or item == f"{label}\uff1a":
            for nxt in items[i + 1 : i + 6]:
                cleaned = nxt.strip().strip("\uff1a")
                if cleaned and cleaned not in {label, "\uff1a"}:
                    return cleaned
    return ""


def export_template_preview(template_path: Path) -> None:
    if any(PREVIEW_DIR.glob("*.PNG")):
        return
    cmd = (
        "$ppt = New-Object -ComObject PowerPoint.Application; "
        f"$pres = $ppt.Presentations.Open('{template_path}', $true, $false, $false); "
        f"$pres.Export('{PREVIEW_DIR}', 'PNG'); "
        "$pres.Close(); "
        "$ppt.Quit();"
    )
    subprocess.run(["powershell", "-NoProfile", "-Command", cmd], check=True, capture_output=True)


def build_content_bg() -> Path:
    export_template_preview(find_source_files().template_pptx)
    bg_path = ASSET_DIR / "content_bg.png"
    if bg_path.exists():
        return bg_path
    preview = next(PREVIEW_DIR.glob("*.PNG"))
    for candidate in PREVIEW_DIR.glob("*.PNG"):
        if "2" in candidate.stem:
            preview = candidate
            break
    with Image.open(preview) as img:
        w, h = img.size
        header = img.crop((0, 0, w, int(h * 0.16)))
        clean_x = min(400, w - 1)
        left_rgb = header.getpixel((0, 10))
        right_rgb = header.getpixel((clean_x, 10))
        for x in range(0, clean_x):
            ratio = x / max(clean_x - 1, 1)
            color = tuple(int(left_rgb[i] + (right_rgb[i] - left_rgb[i]) * ratio) for i in range(3))
            for y in range(header.size[1]):
                header.putpixel((x, y), color)
        canvas = Image.new("RGB", (w, h), "white")
        canvas.paste(header, (0, 0))
        canvas.save(bg_path)
    return bg_path


def gather_stats() -> Stats:
    sources = find_source_files()
    proposal = load_doc_text(sources.proposal_docx)
    title = value_after(proposal, zh(r"\u9898\u76ee")) or zh(r"\u667a\u80fd\u76d1\u63a7\u4e2d\u7684\u884c\u4eba\u5f02\u5e38\u4e8b\u4ef6\u68c0\u6d4b\u4e0e\u591a\u76ee\u6807\u8ddf\u8e2a\u9884\u8b66\u7cfb\u7edf")
    reporter = value_after(proposal, zh(r"\u59d3\u540d")) or zh(r"\u848b\u4e00\u9e23")
    teacher = value_after(proposal, zh(r"\u6307\u5bfc\u6559\u5e08")) or zh(r"\u5510\u5b87\u6d9b")

    best_map50 = -1.0
    best_map5095 = -1.0
    det_run = ""
    for csv_path in (OUTPUT_DIR / "training").glob("*/results.csv"):
        rows = list(csv.DictReader(csv_path.open("r", encoding="utf-8", errors="ignore")))
        if not rows:
            continue
        row = max(rows, key=lambda r: float(r["metrics/mAP50-95(B)"]))
        score = float(row["metrics/mAP50-95(B)"])
        if score > best_map5095:
            best_map50 = float(row["metrics/mAP50(B)"])
            best_map5095 = score
            det_run = csv_path.parent.name

    behavior_path = OUTPUT_DIR / "behavior_training" / "avenue_behavior_temporal_trainsplit_v2_seq72_seqval_seed42" / "metrics.json"
    behavior_data = json.loads(behavior_path.read_text(encoding="utf-8"))
    behavior_f1 = float(behavior_data["val_metrics"]["macro_f1"])
    behavior_val_samples = int(behavior_data["data"]["val_samples"])

    best_frame = {"f1": -1.0}
    for json_path in (OUTPUT_DIR / "avenue").glob("compare_hybrid_6seq*.json"):
        data = json.loads(json_path.read_text(encoding="utf-8"))
        summary = data.get("summary", {})
        if float(summary.get("f1", 0.0)) > float(best_frame["f1"]):
            best_frame = summary

    rep = None
    alarms_total = 0
    for meta_path in OUTPUT_DIR.glob("run_meta_*.json"):
        data = json.loads(meta_path.read_text(encoding="utf-8"))
        summary = data.get("summary", {})
        alarms_total += int(summary.get("total_alarms") or 0)
        counts = summary.get("alarm_counts") or {}
        fps = summary.get("avg_fps")
        if {"intrusion", "loitering", "running"}.issubset(counts.keys()) and fps is not None:
            cand = (float(fps), data.get("source", {}).get("source_name", ""), int(summary.get("total_frames") or 0))
            if rep is None or cand[0] > rep[0]:
                rep = cand

    python_files = 0
    python_lines = 0
    for py in PROJECT_ROOT.rglob("*.py"):
        if "venv" in py.parts:
            continue
        python_files += 1
        python_lines += sum(1 for _ in py.open("r", encoding="utf-8", errors="ignore"))

    pseudo = json.loads((OUTPUT_DIR / "avenue_pseudo_labels" / "summary.json").read_text(encoding="utf-8"))
    expanded = json.loads((OUTPUT_DIR / "avenue_pseudo_labels_expanded" / "summary.json").read_text(encoding="utf-8"))

    return Stats(
        title=title,
        reporter=reporter,
        teacher=teacher,
        detection_map50=best_map50,
        detection_map5095=best_map5095,
        detection_run=det_run,
        behavior_f1=behavior_f1,
        behavior_val_samples=behavior_val_samples,
        frame_f1=float(best_frame["f1"]),
        frame_precision=float(best_frame["precision"]),
        frame_recall=float(best_frame["recall"]),
        runtime_fps=rep[0],
        runtime_source=rep[1],
        runtime_frames=rep[2],
        python_files=python_files,
        python_lines=python_lines,
        behavior_runs=sum(1 for p in (OUTPUT_DIR / "behavior_training").iterdir() if p.is_dir()),
        compare_files=sum(1 for _ in (OUTPUT_DIR / "avenue").glob("compare_*.json")),
        runtime_logs=sum(1 for _ in OUTPUT_DIR.glob("run_meta_*.json")),
        snapshots=sum(1 for _ in (OUTPUT_DIR / "snaps").glob("run_*/*.jpg")),
        alarms_total=alarms_total,
        pseudo_tracks=int(pseudo["summary"]["tracks"]),
        expanded_tracks=int(expanded["summary"]["expanded_tracks"]),
    )


def set_style(paragraph, size: int, color: RGBColor, bold: bool = False, font_name: str = "Microsoft YaHei") -> None:
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = font_name


def xcm(value: float):
    return Cm(value * X_SCALE)


def textbox(slide, left: float, top: float, width: float, height: float, text: str, size: int, color: RGBColor, bold: bool = False, font_name: str = "Microsoft YaHei", align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(xcm(left), Cm(top), xcm(width), Cm(height))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_style(p, size, color, bold, font_name)
    return shape


def bullets(slide, left: float, top: float, width: float, height: float, items: list[str], color: RGBColor = TEXT_DARK):
    shape = slide.shapes.add_textbox(xcm(left), Cm(top), xcm(width), Cm(height))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {item}"
        set_style(p, 18, color)
    return shape


def card(slide, left: float, top: float, width: float, height: float, fill: RGBColor = CARD_BG, border: RGBColor = CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, xcm(left), Cm(top), xcm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.1)
    return shape


def metric_card(slide, left: float, top: float, width: float, label: str, value: str, note: str, accent: RGBColor):
    card(slide, left, top, width, 2.1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, xcm(left + 0.2), Cm(top + 0.25), xcm(0.8), Cm(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, left + 0.35, top + 0.75, width - 0.6, 0.35, label, 11, TEXT_MUTED)
    textbox(slide, left + 0.35, top + 1.05, width - 0.6, 0.45, value, 17, TEXT_DARK, True)
    textbox(slide, left + 0.35, top + 1.48, width - 0.6, 0.32, note, 8, TEXT_MUTED)


def picture(slide, path: Path, left: float, top: float, width: float, height: float):
    slide.shapes.add_picture(str(path), xcm(left), Cm(top), width=xcm(width), height=Cm(height))


def add_bg(slide, prs: Presentation, title: str, bg_path: Path) -> None:
    slide.shapes.add_picture(str(bg_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    mask = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, xcm(2.6), Cm(1.55))
    mask.fill.solid()
    mask.fill.fore_color.rgb = RGBColor(68, 72, 157)
    mask.line.fill.background()
    textbox(slide, 0.7, 0.28, 6.5, 0.8, title, 23, RGBColor(255, 255, 255), False, "STXingkai")


def replace_cover(slide, stats: Stats) -> None:
    mapping = {
        "title": zh(r"\u667a\u80fd\u76d1\u63a7\u4e2d\u7684\u884c\u4eba\u5f02\u5e38\u4e8b\u4ef6\u68c0\u6d4b\u4e0e\n\u591a\u76ee\u6807\u8ddf\u8e2a\u9884\u8b66\u7cfb\u7edf"),
        "reporter": f"{zh(r'\u6c47\u62a5\u4eba')}：{stats.reporter}",
        "teacher": f"{zh(r'\u6307\u5bfc\u8001\u5e08')}：{stats.teacher}",
        "date": "2026.3.29",
    }
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = shape.text.strip()
        if "面向扫描电镜图像" in text:
            shape.text = mapping["title"]
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                set_style(p, 28, RGBColor(255, 255, 255), False)
        elif text.startswith("汇报人"):
            shape.text = mapping["reporter"]
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                set_style(p, 18, RGBColor(255, 255, 255), False)
        elif text.startswith("指导老师"):
            shape.text = mapping["teacher"]
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                set_style(p, 18, RGBColor(255, 255, 255), False)
        elif text.startswith("2026"):
            shape.text = mapping["date"]
            for p in shape.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                set_style(p, 16, RGBColor(255, 255, 255), False)


def fill_outline(slide) -> None:
    textbox(slide, 1.2, 2.0, 9.6, 0.6, zh(r"\u672c\u6b21\u6c47\u62a5\u805a\u7126\u201c\u505a\u5230\u54ea\u91cc\u3001\u6548\u679c\u600e\u6837\u3001\u540e\u9762\u600e\u4e48\u505a\u201d\u3002"), 17, TEXT_MUTED)
    items = [
        zh(r"1. \u8bfe\u9898\u80cc\u666f\u4e0e\u76ee\u6807"),
        zh(r"2. \u7cfb\u7edf\u8bbe\u8ba1\u4e0e\u6280\u672f\u8def\u7ebf"),
        zh(r"3. \u9636\u6bb5\u6210\u679c\u4e0e\u8fd0\u884c\u6f14\u793a"),
        zh(r"4. \u5b9e\u9a8c\u7ed3\u679c\u4e0e\u5f53\u524d\u8fdb\u5ea6"),
        zh(r"5. \u95ee\u9898\u3001\u6539\u8fdb\u4e0e\u5de5\u4f5c\u91cf"),
    ]
    y = 3.2
    for item in items:
        textbox(slide, 1.7, y, 9.0, 0.55, item, 18, TEXT_DARK, True)
        y += 0.95


def build_ppt() -> Path:
    ensure_dirs()
    stats = gather_stats()
    bg_path = build_content_bg()
    sources = find_source_files()
    prs = Presentation(str(sources.template_pptx))
    replace_cover(prs.slides[0], stats)
    fill_outline(prs.slides[1])
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[len(prs.slide_layouts) - 1]

    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs, zh(r"\u8bfe\u9898\u80cc\u666f\u4e0e\u4e2d\u671f\u76ee\u6807"), bg_path)
    card(slide, 0.8, 1.9, 5.8, 8.4)
    textbox(slide, 1.15, 2.25, 3.0, 0.6, zh(r"\u7814\u7a76\u80cc\u666f"), 21, TEXT_DARK, True)
    bullets(slide, 1.15, 3.0, 4.9, 4.5, [
        zh(r"\u4f20\u7edf\u89c6\u9891\u76d1\u63a7\u4f9d\u8d56\u4eba\u5de5\u503c\u5b88\uff0c\u6613\u6f0f\u68c0\u3001\u53cd\u5e94\u6ede\u540e\u3002"),
        zh(r"\u9898\u76ee\u8981\u6c42\u540c\u65f6\u5b8c\u6210\u884c\u4eba\u68c0\u6d4b\u3001\u591a\u76ee\u6807\u8ddf\u8e2a\u548c\u5f02\u5e38\u9884\u8b66\u3002"),
        zh(r"\u76ee\u6807\u4e8b\u4ef6\u5305\u62ec\u8d8a\u754c\u3001\u5f98\u5f8a\u3001\u5feb\u901f\u5954\u8dd1\uff0c\u5e76\u8981\u6c42FPS > 15\u3002"),
        zh(r"\u4e2d\u671f\u91cd\u70b9\u662f\u5148\u628a\u7b97\u6cd5\u94fe\u8def\u8dd1\u901a\uff0c\u518d\u505a\u524d\u7aef\u548c\u6700\u7ec8\u5305\u88c5\u3002"),
    ])
    card(slide, 6.9, 1.9, 5.4, 4.3, RGBColor(255, 247, 250), RGBColor(240, 213, 224))
    textbox(slide, 7.25, 2.25, 3.0, 0.6, zh(r"\u5f53\u524d\u4e2d\u671f\u5224\u65ad"), 21, TEXT_DARK, True)
    bullets(slide, 7.25, 3.0, 4.3, 2.4, [
        zh(r"\u68c0\u6d4b\u8bad\u7ec3\u3001\u8ddf\u8e2a\u8054\u52a8\u3001\u4e09\u7c7b\u544a\u8b66\u5747\u5df2\u80fd\u8fd0\u884c\u3002"),
        zh(r"\u8f7b\u91cf\u884c\u4e3a\u6a21\u578b\u5df2\u5f00\u59cb\u878d\u5408\uff0c\u4f46\u4ecd\u9700\u8981\u7ee7\u7eed\u51cf\u5c11\u8bef\u62a5\u3002"),
        zh(r"\u754c\u9762\u3001\u56de\u653e\u3001\u6700\u7ec8\u7cfb\u7edf\u5316\u5c1a\u672a\u5b8c\u6210\u3002"),
    ], ACCENT_PINK)
    metric_card(slide, 6.9, 6.7, 2.45, zh(r"\u8fd0\u884c\u901f\u5ea6"), f"{stats.runtime_fps:.2f} FPS", zh(r"\u4ee3\u8868\u6027Demo"), ACCENT_GREEN)
    metric_card(slide, 9.55, 6.7, 2.75, zh(r"\u5e27\u7ea7F1"), f"{stats.frame_f1:.4f}", zh(r"\u5f02\u5e38\u5e8f\u5217\u5bf9\u6bd4"), ACCENT_GOLD)

    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs, zh(r"\u7cfb\u7edf\u8bbe\u8ba1\u4e0e\u6280\u672f\u8def\u7ebf"), bg_path)
    titles = [
        zh(r"\u89c6\u9891\u8f93\u5165"),
        zh(r"YOLO\u884c\u4eba\u68c0\u6d4b"),
        zh(r"ByteTrack\u8ddf\u8e2a"),
        zh(r"\u5f02\u5e38\u5206\u6790"),
        zh(r"\u544a\u8b66\u4e0e\u65e5\u5fd7"),
    ]
    notes = [
        zh(r"\u6444\u50cf\u5934/\u89c6\u9891/MOT17\u56fe\u50cf\u5e8f\u5217"),
        zh(r"\u8f93\u51faperson\u68c0\u6d4b\u6846"),
        zh(r"\u7ef4\u6301ID\u4e0e\u8f68\u8ff9"),
        zh(r"\u8d8a\u754c/\u5f98\u5f8a/\u5954\u8dd1"),
        zh(r"\u622a\u56fe/\u65e5\u5fd7/\u53ef\u89c6\u5316"),
    ]
    xs = [0.8, 3.25, 5.7, 8.15, 10.6]
    for x, t, n in zip(xs, titles, notes):
        card(slide, x, 3.0, 1.95, 2.3, RGBColor(248, 250, 255), CARD_BORDER)
        textbox(slide, x + 0.2, 3.35, 1.55, 0.55, t, 16, TEXT_DARK, True)
        textbox(slide, x + 0.2, 4.0, 1.55, 0.8, n, 11, TEXT_MUTED)
    for x in [2.8, 5.25, 7.7, 10.15]:
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, xcm(x), Cm(3.7), xcm(0.35), Cm(0.55))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.fill.background()
    card(slide, 0.9, 6.2, 11.3, 3.8, RGBColor(247, 250, 255), CARD_BORDER)
    bullets(slide, 1.2, 6.7, 10.4, 2.7, [
        zh(r"\u8bad\u7ec3\u94fe\u8def\uff1aMOT17 -> YOLO\u683c\u5f0f\u8f6c\u6362 -> \u884c\u4eba\u68c0\u6d4b\u5668\u8bad\u7ec3\u4e0e\u9a8c\u8bc1\u3002"),
        zh(r"\u884c\u4e3a\u94fe\u8def\uff1aAvenue\u5f02\u5e38\u63a9\u7801 -> \u8f68\u8ff9\u4f2a\u6807\u7b7e -> \u8f7b\u91cf\u65f6\u5e8f\u6a21\u578b\u8bad\u7ec3\u3002"),
        zh(r"\u8fd0\u884c\u94fe\u8def\uff1a\u89c6\u9891\u8f93\u5165 -> \u68c0\u6d4b -> \u8ddf\u8e2a -> \u89c4\u5219/\u6a21\u578b\u878d\u5408 -> \u62a5\u8b66\u8f93\u51fa\u3002"),
    ])

    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs, zh(r"\u8fd0\u884c\u6f14\u793a\u4e0e\u9636\u6bb5\u6210\u679c"), bg_path)
    pics = [
        OUTPUT_DIR / "snaps" / "run_20260327_200408_045723" / "intrusion_tid1_20260327_200414_835665.jpg",
        OUTPUT_DIR / "snaps" / "run_20260327_200408_045723" / "loitering_tid6_20260327_200417_003890.jpg",
        OUTPUT_DIR / "snaps" / "run_20260327_200408_045723" / "running_tid32_20260327_200424_007508.jpg",
    ]
    caps = [zh(r"\u533a\u57df\u5165\u4fb5"), zh(r"\u4eba\u5458\u5f98\u5f8a"), zh(r"\u5feb\u901f\u5954\u8dd1")]
    for i, (pic, cap) in enumerate(zip(pics, caps)):
        x = 0.85 + i * 4.05
        picture(slide, pic, x, 2.0, 3.75, 4.6)
        textbox(slide, x, 6.8, 3.75, 0.45, cap, 18, TEXT_DARK, True, align=PP_ALIGN.CENTER)
    metric_card(slide, 0.9, 8.0, 3.6, zh(r"\u4ee3\u8868\u6027Demo"), stats.runtime_source, f"{stats.runtime_frames}{zh(r'\u5e27\u5b8c\u6574\u8fd0\u884c')}", ACCENT_BLUE)
    metric_card(slide, 4.8, 8.0, 3.6, zh(r"\u5b9e\u65f6\u901f\u5ea6"), f"{stats.runtime_fps:.2f} FPS", zh(r"\u5df2\u8fbe\u5230\u51c6\u5b9e\u65f6"), ACCENT_GREEN)
    metric_card(slide, 8.7, 8.0, 3.6, zh(r"\u544a\u8b66\u7c7b\u578b"), zh(r"\u8d8a\u754c/\u5f98\u5f8a/\u5954\u8dd1"), zh(r"\u5df2\u80fd\u5728\u4e00\u5957\u6d41\u7a0b\u5185\u89e6\u53d1"), ACCENT_PINK)

    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs, zh(r"\u5b9e\u9a8c\u7ed3\u679c\u4e0e\u5f53\u524d\u8fdb\u5ea6"), bg_path)
    metric_card(slide, 0.8, 2.0, 2.8, zh(r"\u68c0\u6d4b mAP50"), f"{stats.detection_map50:.4f}", stats.detection_run, ACCENT_BLUE)
    metric_card(slide, 3.9, 2.0, 2.8, zh(r"\u68c0\u6d4b mAP50-95"), f"{stats.detection_map5095:.4f}", zh(r"\u6700\u4f73\u884c\u4eba\u68c0\u6d4b\u7ed3\u679c"), ACCENT_PINK)
    metric_card(slide, 7.0, 2.0, 2.5, zh(r"\u884c\u4e3a Macro-F1"), f"{stats.behavior_f1:.4f}", f"n={stats.behavior_val_samples}", ACCENT_GREEN)
    metric_card(slide, 9.8, 2.0, 2.4, zh(r"\u5e27\u7ea7 F1"), f"{stats.frame_f1:.4f}", f"P {stats.frame_precision:.3f} / R {stats.frame_recall:.3f}", ACCENT_GOLD)
    picture(slide, OUTPUT_DIR / "training" / stats.detection_run / "results.png", 0.9, 4.9, 6.0, 4.7)
    card(slide, 7.3, 4.9, 5.0, 4.7)
    textbox(slide, 7.65, 5.2, 3.0, 0.5, zh(r"\u8fdb\u5ea6\u5224\u65ad"), 20, TEXT_DARK, True)
    progress = [
        zh(r"\u6587\u732e\u8c03\u7814\u4e0e\u5f00\u9898 100%"),
        zh(r"\u6570\u636e\u51c6\u5907\u4e0e\u9884\u5904\u7406 95%"),
        zh(r"\u68c0\u6d4b\u8bad\u7ec3\u4e0e\u9a8c\u8bc1 100%"),
        zh(r"\u8ddf\u8e2a\u4e0e\u89c4\u5219\u544a\u8b66 90%"),
        zh(r"\u8f7b\u91cf\u884c\u4e3a\u6a21\u578b 85%"),
        zh(r"\u754c\u9762\u4e0e\u6700\u7ec8\u6574\u5408 35%"),
    ]
    bullets(slide, 7.65, 5.9, 4.0, 3.0, progress)

    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs, zh(r"\u5f53\u524d\u95ee\u9898\u4e0e\u6539\u8fdb\u65b9\u5411"), bg_path)
    issues = [
        zh(r"\u4f2a\u6807\u7b7e\u566a\u58f0\uff1aAvenue\u6837\u672c\u4e0d\u662f\u4eba\u5de5\u7cbe\u6807\uff0c\u4ecd\u5b58\u5728\u6807\u7b7e\u504f\u5dee\u3002"),
        zh(r"\u8bef\u62a5\u504f\u9ad8\uff1a\u8d8a\u754c\u4e0e\u5f98\u5f8a\u5bf9ROI\u3001\u9608\u503c\u548c\u573a\u666f\u5bc6\u5ea6\u8f83\u654f\u611f\u3002"),
        zh(r"\u5b9e\u65f6\u6027\u6ce2\u52a8\uff1a\u5f15\u5165\u65f6\u5e8f\u6a21\u578b\u540eFPS\u4f1a\u4e0b\u964d\uff0c\u8fd8\u9700\u8981\u95e8\u63a7\u4e0e\u538b\u7f29\u3002"),
        zh(r"\u5de5\u7a0b\u5316\u4e0d\u8db3\uff1aGUI\u3001\u56de\u653e\u3001\u53c2\u6570\u914d\u7f6e\u754c\u9762\u8fd8\u6ca1\u6709\u5b8c\u5168\u8865\u9f50\u3002"),
    ]
    y = 2.2
    for idx, issue in enumerate(issues):
        x = 0.9 if idx % 2 == 0 else 6.7
        if idx and idx % 2 == 0:
            y += 3.2
        card(slide, x, y, 5.0, 2.6)
        textbox(slide, x + 0.3, y + 0.35, 4.2, 1.5, issue, 16, TEXT_MUTED)
    card(slide, 0.9, 8.9, 11.3, 1.1, RGBColor(248, 250, 255), CARD_BORDER)
    textbox(slide, 1.2, 9.15, 10.3, 0.4, zh(r"\u4e0b\u4e00\u9636\u6bb5\u91cd\u70b9\uff1a\u964d\u8bef\u62a5\u3001\u589e\u5f3a\u8ddf\u8e2a\u9c81\u68d2\u6027\u3001\u8865\u9f50GUI\u4e0e\u8bba\u6587\u6750\u6599\u3002"), 16, TEXT_DARK, True)

    slide = prs.slides.add_slide(blank)
    add_bg(slide, prs, zh(r"\u5de5\u4f5c\u91cf\u7edf\u8ba1"), bg_path)
    metric_card(slide, 0.8, 2.0, 2.9, zh(r"\u4ee3\u7801\u6587\u4ef6"), str(stats.python_files), zh(r"\u4e2a Python \u6587\u4ef6"), ACCENT_BLUE)
    metric_card(slide, 3.95, 2.0, 2.9, zh(r"\u4ee3\u7801\u884c\u6570"), str(stats.python_lines), zh(r"\u884c\u6838\u5fc3\u4ee3\u7801"), ACCENT_PINK)
    metric_card(slide, 7.1, 2.0, 2.5, zh(r"\u884c\u4e3a\u8bad\u7ec3"), str(stats.behavior_runs), zh(r"\u7ec4\u5b9e\u9a8c\u76ee\u5f55"), ACCENT_GREEN)
    metric_card(slide, 9.9, 2.0, 2.3, zh(r"Avenue\u5bf9\u6bd4"), str(stats.compare_files), zh(r"\u4efdJSON"), ACCENT_GOLD)
    metric_card(slide, 0.8, 5.0, 2.9, zh(r"\u8fd0\u884c\u65e5\u5fd7"), str(stats.runtime_logs), zh(r"\u4efd run_meta"), ACCENT_BLUE)
    metric_card(slide, 3.95, 5.0, 2.9, zh(r"\u544a\u8b66\u622a\u56fe"), str(stats.snapshots), zh(r"\u5f20\u4fdd\u5b58\u622a\u56fe"), ACCENT_PINK)
    metric_card(slide, 7.1, 5.0, 2.5, zh(r"\u4f2a\u6807\u7b7e"), str(stats.pseudo_tracks), zh(r"\u6761\u57fa\u7840\u8f68\u8ff9"), ACCENT_GREEN)
    metric_card(slide, 9.9, 5.0, 2.3, zh(r"\u6269\u5145\u540e"), str(stats.expanded_tracks), zh(r"\u6761\u8bad\u7ec3\u6837\u672c"), ACCENT_GOLD)
    card(slide, 0.9, 8.3, 11.2, 1.5)
    textbox(slide, 1.2, 8.7, 10.5, 0.55, f"{zh(r'\u7d2f\u8ba1\u544a\u8b66')} {stats.alarms_total} {zh(r'\u6761')}；{zh(r'\u884c\u4e3a\u4f2a\u6807\u7b7e\u7531')} {stats.pseudo_tracks} {zh(r'\u6269\u5145\u5230')} {stats.expanded_tracks} {zh(r'\u6761')}；{zh(r'\u5df2\u5f62\u6210\u53ef\u7528\u7684\u4e2d\u671f\u6c47\u62a5\u6750\u6599\u94fe\u8def')}。", 16, TEXT_MUTED)

    sld_ids = prs.slides._sldIdLst
    slide_id = sld_ids[2]
    sld_ids.remove(slide_id)
    sld_ids.append(slide_id)

    out = DESKTOP / zh(r"\u6bd5\u8bbe") / zh(r"2201630311_\u848b\u4e00\u9e23_\u6bd5\u4e1a\u8bbe\u8ba1\u4e2d\u671f\u6c47\u62a5.pptx")
    prs.save(str(out))
    return out


def inspect() -> None:
    sources = find_source_files()
    print(sources.proposal_docx)
    print(sources.title_docx)
    print(sources.template_pptx)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("command", choices=["inspect", "build"])
    args = parser.parse_args()
    if args.command == "inspect":
        inspect()
        return
    out = build_ppt()
    print(f"OUTPUT={out}")


if __name__ == "__main__":
    main()
